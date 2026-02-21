#!/usr/bin/env python3
"""VideoCut 説明書 PowerPoint 生成スクリプト"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# --- カラーテーマ ---
BG_DARK = RGBColor(0x14, 0x14, 0x14)
BG_SURFACE = RGBColor(0x1E, 0x1E, 0x1E)
BG_CARD = RGBColor(0x25, 0x25, 0x25)
ACCENT = RGBColor(0x4F, 0xC3, 0xF7)
GREEN = RGBColor(0x4C, 0xAF, 0x50)
RED = RGBColor(0xF4, 0x43, 0x36)
YELLOW = RGBColor(0xFF, 0xB7, 0x4D)
WHITE = RGBColor(0xE0, 0xE0, 0xE0)
GRAY = RGBColor(0x99, 0x99, 0x99)
DARK_GRAY = RGBColor(0x55, 0x55, 0x55)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

OUT = os.path.join(os.path.dirname(__file__), "VideoCut_Manual.pptx")


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, corner_radius=None):
    if corner_radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.adjustments[0] = corner_radius
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE,
             bold=False, alignment=PP_ALIGN.LEFT, font_name="Helvetica Neue"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=WHITE, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing

        if isinstance(item, tuple):
            # (label, description) 形式
            run1 = p.add_run()
            run1.text = item[0]
            run1.font.size = Pt(font_size)
            run1.font.color.rgb = ACCENT
            run1.font.bold = True
            run1.font.name = "Helvetica Neue"
            run2 = p.add_run()
            run2.text = "  " + item[1]
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = color
            run2.font.name = "Helvetica Neue"
        else:
            p.text = item
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.name = "Helvetica Neue"
        p.level = 0
    return txBox


def add_key_badge(slide, left, top, key_text, width=None):
    """キーボードキーのバッジを描画"""
    w = width or Inches(0.7)
    h = Inches(0.38)
    shape = add_rect(slide, left, top, w, h, RGBColor(0x33, 0x33, 0x33), corner_radius=0.15)
    shape.line.color.rgb = RGBColor(0x55, 0x55, 0x55)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = key_text
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    p.font.name = "Menlo"
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    return shape


# ============================================================
# スライド作成
# ============================================================
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # --- 1. タイトルスライド ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_DARK)

    # 中央のカード
    card_w, card_h = Inches(8), Inches(4)
    card_x = (SLIDE_W - card_w) // 2
    card_y = (SLIDE_H - card_h) // 2 - Inches(0.3)
    add_rect(slide, card_x, card_y, card_w, card_h, BG_SURFACE, corner_radius=0.05)

    # アクセントライン
    add_rect(slide, card_x, card_y, card_w, Inches(0.06), ACCENT)

    add_text(slide, card_x, card_y + Inches(0.5), card_w, Inches(1),
             "VideoCut", font_size=52, color=ACCENT, bold=True,
             alignment=PP_ALIGN.CENTER)
    add_text(slide, card_x, card_y + Inches(1.6), card_w, Inches(0.6),
             "シンプル動画カットエディタ for macOS", font_size=22, color=WHITE,
             alignment=PP_ALIGN.CENTER)
    add_text(slide, card_x, card_y + Inches(2.4), card_w, Inches(0.5),
             "動画をプレビューしながら、必要な部分だけをカット＆結合",
             font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

    add_text(slide, card_x, card_y + Inches(3.2), card_w, Inches(0.5),
             "Version 1.0", font_size=14, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    # --- 2. 概要 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "VideoCut とは", font_size=32, color=ACCENT, bold=True)
    add_rect(slide, Inches(0.8), Inches(1.05), Inches(2), Inches(0.04), ACCENT)

    add_text(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.8),
             "macOS で動作するシンプルな動画カット編集ツールです。\n"
             "ブラウザベースのUIで、動画のプレビュー・カット・結合・上下反転をサポートします。",
             font_size=18, color=WHITE)

    # 特徴カード
    features = [
        ("プレビュー再生", "動画を再生しながら\nカット位置を確認", ACCENT),
        ("IN / OUT 編集", "残したい区間を\n自由にマーク", GREEN),
        ("無劣化エクスポート", "ffmpeg による\n高速・無劣化出力", YELLOW),
        ("上下反転", "プレビュー＆出力\n両方に反映", RED),
    ]
    card_w_each = Inches(2.7)
    gap = Inches(0.3)
    start_x = Inches(0.8)
    for i, (title, desc, color) in enumerate(features):
        x = start_x + i * (card_w_each + gap)
        y = Inches(3.0)
        add_rect(slide, x, y, card_w_each, Inches(2.6), BG_CARD, corner_radius=0.04)
        add_rect(slide, x, y, card_w_each, Inches(0.06), color)
        add_text(slide, x + Inches(0.3), y + Inches(0.4), card_w_each - Inches(0.6), Inches(0.5),
                 title, font_size=18, color=color, bold=True)
        add_text(slide, x + Inches(0.3), y + Inches(1.1), card_w_each - Inches(0.6), Inches(1.2),
                 desc, font_size=15, color=GRAY)

    # --- 3. インストール ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "インストール方法", font_size=32, color=ACCENT, bold=True)
    add_rect(slide, Inches(0.8), Inches(1.05), Inches(2.8), Inches(0.04), ACCENT)

    steps = [
        ("Step 1", "必要環境の確認", "macOS 12 以降\nffmpeg がインストール済みであること"),
        ("Step 2", "ffmpeg のインストール", "ターミナルで以下を実行:\n  brew install ffmpeg"),
        ("Step 3", "アプリの配置", "VideoCut.dmg を開き、\nVideoCut.app を アプリケーション フォルダにコピー"),
        ("Step 4", "初回起動", "右クリック →「開く」を選択\n（初回は開発元確認ダイアログが出ます）"),
    ]
    for i, (step, title, desc) in enumerate(steps):
        y = Inches(1.6) + i * Inches(1.35)
        # ステップ番号バッジ
        add_rect(slide, Inches(0.8), y, Inches(1.2), Inches(1.05), BG_CARD, corner_radius=0.06)
        add_text(slide, Inches(0.8), y + Inches(0.15), Inches(1.2), Inches(0.4),
                 step, font_size=14, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
        add_text(slide, Inches(0.8), y + Inches(0.5), Inches(1.2), Inches(0.4),
                 "●", font_size=20, color=ACCENT, alignment=PP_ALIGN.CENTER)
        # 内容
        add_text(slide, Inches(2.3), y + Inches(0.05), Inches(5), Inches(0.4),
                 title, font_size=18, color=WHITE, bold=True)
        add_text(slide, Inches(2.3), y + Inches(0.45), Inches(9), Inches(0.6),
                 desc, font_size=14, color=GRAY)

    # --- 4. 画面構成 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "画面構成", font_size=32, color=ACCENT, bold=True)
    add_rect(slide, Inches(0.8), Inches(1.05), Inches(1.8), Inches(0.04), ACCENT)

    # UIモックアップ
    mock_x, mock_y = Inches(1.5), Inches(1.5)
    mock_w, mock_h = Inches(10), Inches(5.5)

    # 全体枠
    add_rect(slide, mock_x, mock_y, mock_w, mock_h, BG_SURFACE, corner_radius=0.02)

    # ヘッダー
    add_rect(slide, mock_x, mock_y, mock_w, Inches(0.55), BG_CARD)
    add_text(slide, mock_x + Inches(0.2), mock_y + Inches(0.08), Inches(1.2), Inches(0.4),
             "VideoCut", font_size=14, color=ACCENT, bold=True)
    add_text(slide, mock_x + Inches(1.8), mock_y + Inches(0.1), Inches(4), Inches(0.35),
             "sample.mp4 (1920x1080, 30fps)", font_size=11, color=GRAY)

    # ヘッダーボタン
    for j, (label, c) in enumerate([("開く", WHITE), ("上下反転", WHITE), ("エクスポート", GREEN)]):
        bx = mock_x + mock_w - Inches(1.1) * (3 - j) - Inches(0.2)
        shape = add_rect(slide, bx, mock_y + Inches(0.1), Inches(1.0), Inches(0.35),
                         RGBColor(0x30, 0x30, 0x30), corner_radius=0.12)
        shape.line.color.rgb = RGBColor(0x44, 0x44, 0x44)
        shape.line.width = Pt(1)
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(10)
        p.font.color.rgb = c
        p.alignment = PP_ALIGN.CENTER

    # プレビューエリア
    preview_y = mock_y + Inches(0.55)
    add_rect(slide, mock_x, preview_y, mock_w, Inches(2.8), RGBColor(0x05, 0x05, 0x05))
    add_text(slide, mock_x, preview_y + Inches(1.0), mock_w, Inches(0.5),
             "[ 動画プレビュー ]", font_size=20, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    # タイムライン
    tl_y = preview_y + Inches(2.8)
    add_rect(slide, mock_x, tl_y, mock_w, Inches(0.8), BG_SURFACE)
    add_text(slide, mock_x + Inches(0.2), tl_y + Inches(0.05), Inches(2), Inches(0.3),
             "00:01:23.45", font_size=13, color=ACCENT, bold=True, font_name="Menlo")
    add_text(slide, mock_x + Inches(2.5), tl_y + Inches(0.08), Inches(2), Inches(0.3),
             "/ 00:05:30.00", font_size=11, color=GRAY, font_name="Menlo")
    # トラック
    track_y = tl_y + Inches(0.35)
    add_rect(slide, mock_x + Inches(0.2), track_y, mock_w - Inches(0.4), Inches(0.25),
             RGBColor(0x22, 0x22, 0x22), corner_radius=0.1)
    # セグメント例
    add_rect(slide, mock_x + Inches(1.5), track_y, Inches(2.0), Inches(0.25),
             RGBColor(0x2E, 0x5C, 0x3A))
    # 再生ヘッド
    add_rect(slide, mock_x + Inches(3.0), track_y - Inches(0.05), Inches(0.03), Inches(0.35), ACCENT)

    # コントロール
    ctrl_y = tl_y + Inches(0.8)
    add_rect(slide, mock_x, ctrl_y, mock_w, Inches(0.5), BG_SURFACE)
    ctrl_labels = ["⏮", "⏪", "◀", " ▶ ", "▶", "⏩", "⏭", "", "I-IN", "O-OUT", "全選択", "＋追加"]
    cx = mock_x + Inches(1.5)
    for label in ctrl_labels:
        if label == "":
            cx += Inches(0.15)
            continue
        w = Inches(0.65) if len(label) > 2 else Inches(0.45)
        shape = add_rect(slide, cx, ctrl_y + Inches(0.08), w, Inches(0.32),
                         RGBColor(0x30, 0x30, 0x30), corner_radius=0.12)
        shape.line.color.rgb = RGBColor(0x44, 0x44, 0x44)
        shape.line.width = Pt(1)
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(9)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        cx += w + Inches(0.06)

    # カットリスト
    cl_y = ctrl_y + Inches(0.5)
    add_rect(slide, mock_x, cl_y, mock_w, Inches(0.85), BG_SURFACE)
    add_text(slide, mock_x + Inches(0.2), cl_y + Inches(0.05), Inches(4), Inches(0.3),
             "カットリスト（保持するセグメント）", font_size=10, color=GRAY)
    add_text(slide, mock_x + Inches(0.3), cl_y + Inches(0.35), Inches(8), Inches(0.2),
             "1.  00:00:10.00 → 00:00:30.00  (00:00:20.00)", font_size=10,
             color=WHITE, font_name="Menlo")
    add_text(slide, mock_x + Inches(0.3), cl_y + Inches(0.58), Inches(8), Inches(0.2),
             "2.  00:01:00.00 → 00:01:45.00  (00:00:45.00)", font_size=10,
             color=WHITE, font_name="Menlo")

    # 注釈ラベル (右側)
    annotations = [
        (mock_y + Inches(0.1), "ヘッダー: ファイル操作・反転・エクスポート"),
        (preview_y + Inches(1.0), "プレビュー: 動画をリアルタイム再生"),
        (tl_y + Inches(0.15), "タイムライン: シーク＆マーカー表示"),
        (ctrl_y + Inches(0.05), "コントロール: 再生操作・IN/OUT設定"),
        (cl_y + Inches(0.15), "カットリスト: 保持セグメント一覧"),
    ]
    for ay, label in annotations:
        # 吹き出し線は省略し、テキストのみ右端に配置
        pass  # レイアウトがモック内に収まるので注釈はモック内のラベルで十分

    # --- 5. 基本操作フロー ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "基本操作フロー", font_size=32, color=ACCENT, bold=True)
    add_rect(slide, Inches(0.8), Inches(1.05), Inches(2.4), Inches(0.04), ACCENT)

    flow_steps = [
        ("1", "動画を開く", "「開く」ボタンをクリック\nファイル選択ダイアログで動画を選択", ACCENT),
        ("2", "プレビュー確認", "再生・シークで内容を確認\nタイムラインをドラッグして移動", WHITE),
        ("3", "IN ポイント設定", "残したい区間の開始位置で\nI キーまたは「I - IN」ボタン", GREEN),
        ("4", "OUT ポイント設定", "残したい区間の終了位置で\nO キーまたは「O - OUT」ボタン", RED),
        ("5", "カットリストに追加", "Enter キーまたは「＋追加」ボタン\n複数セグメント登録可能", YELLOW),
        ("6", "エクスポート", "「エクスポート」→ 確認 → 出力\n元ファイルと同じフォルダに保存", GREEN),
    ]

    for i, (num, title, desc, color) in enumerate(flow_steps):
        col = i % 3
        row = i // 3
        x = Inches(0.8) + col * Inches(4.2)
        y = Inches(1.6) + row * Inches(2.8)

        # カード
        add_rect(slide, x, y, Inches(3.8), Inches(2.3), BG_CARD, corner_radius=0.04)

        # 番号バッジ
        badge = add_rect(slide, x + Inches(0.2), y + Inches(0.2), Inches(0.5), Inches(0.5),
                         color, corner_radius=0.2)
        tf = badge.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = num
        run.font.size = Pt(20)
        run.font.color.rgb = BG_DARK
        run.font.bold = True

        add_text(slide, x + Inches(0.9), y + Inches(0.25), Inches(2.6), Inches(0.4),
                 title, font_size=18, color=WHITE, bold=True)
        add_text(slide, x + Inches(0.3), y + Inches(0.9), Inches(3.2), Inches(1.2),
                 desc, font_size=14, color=GRAY)

        # 矢印（最後以外）
        if i < len(flow_steps) - 1 and col < 2:
            add_text(slide, x + Inches(3.8), y + Inches(0.8), Inches(0.4), Inches(0.5),
                     "→", font_size=24, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    # --- 6. キーボードショートカット ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "キーボードショートカット", font_size=32, color=ACCENT, bold=True)
    add_rect(slide, Inches(0.8), Inches(1.05), Inches(3.8), Inches(0.04), ACCENT)

    shortcuts = [
        ("Space", "再生 / 一時停止"),
        ("← →", "1秒 戻る / 進む"),
        ("Shift + ← →", "10秒 戻る / 進む"),
        (",  .", "1フレーム 戻る / 進む"),
        ("I", "INポイントを現在位置に設定"),
        ("O", "OUTポイントを現在位置に設定"),
        ("A", "全選択（IN=先頭, OUT=末尾）"),
        ("F", "上下反転 ON / OFF"),
        ("Enter", "IN-OUT区間をカットリストに追加"),
    ]

    # 2列レイアウト
    col_items = [shortcuts[:5], shortcuts[5:]]
    for col, items in enumerate(col_items):
        base_x = Inches(0.8) + col * Inches(6.2)
        for i, (key, desc) in enumerate(items):
            y = Inches(1.5) + i * Inches(1.05)
            # カード背景
            add_rect(slide, base_x, y, Inches(5.8), Inches(0.85), BG_CARD, corner_radius=0.04)

            # キーバッジ
            keys = key.split(" + ") if " + " in key else key.split("  ")
            kx = base_x + Inches(0.2)
            for k in keys:
                k = k.strip()
                if not k:
                    continue
                kw = Inches(0.5 + len(k) * 0.12)
                add_key_badge(slide, kx, y + Inches(0.22), k, width=kw)
                kx += kw + Inches(0.1)

            # 説明
            add_text(slide, base_x + Inches(2.2), y + Inches(0.2), Inches(3.4), Inches(0.5),
                     desc, font_size=15, color=WHITE)

    # --- 7. 上下反転機能 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "上下反転機能", font_size=32, color=ACCENT, bold=True)
    add_rect(slide, Inches(0.8), Inches(1.05), Inches(2.2), Inches(0.04), ACCENT)

    add_text(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.6),
             "撮影時に上下逆さまになった動画を正しい向きに修正できます。",
             font_size=18, color=WHITE)

    # 左: OFF状態
    add_rect(slide, Inches(0.8), Inches(2.4), Inches(5.5), Inches(4.2), BG_CARD, corner_radius=0.04)
    add_text(slide, Inches(0.8), Inches(2.6), Inches(5.5), Inches(0.5),
             "反転 OFF", font_size=20, color=GRAY, bold=True, alignment=PP_ALIGN.CENTER)
    add_rect(slide, Inches(1.3), Inches(3.3), Inches(4.5), Inches(2.5), RGBColor(0x05, 0x05, 0x05))
    add_text(slide, Inches(1.3), Inches(4.0), Inches(4.5), Inches(0.8),
             "▲ 通常表示", font_size=18, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    # 右: ON状態
    add_rect(slide, Inches(7.0), Inches(2.4), Inches(5.5), Inches(4.2), BG_CARD, corner_radius=0.04)
    add_rect(slide, Inches(7.0), Inches(2.4), Inches(5.5), Inches(0.06), ACCENT)
    add_text(slide, Inches(7.0), Inches(2.6), Inches(5.5), Inches(0.5),
             "反転 ON", font_size=20, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
    add_rect(slide, Inches(7.5), Inches(3.3), Inches(4.5), Inches(2.5), RGBColor(0x05, 0x05, 0x05))
    add_text(slide, Inches(7.5), Inches(4.0), Inches(4.5), Inches(0.8),
             "▼ 上下反転表示", font_size=18, color=ACCENT, alignment=PP_ALIGN.CENTER)

    # 注意事項
    add_rect(slide, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.7), BG_SURFACE, corner_radius=0.04)
    add_text(slide, Inches(1.1), Inches(6.1), Inches(11), Inches(0.5),
             "注意: 反転ONでエクスポートすると再エンコード（H.264, CRF18）になるため、"
             "反転OFFの無劣化コピーより時間がかかります。",
             font_size=14, color=YELLOW)

    # --- 8. エクスポート ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "エクスポート", font_size=32, color=ACCENT, bold=True)
    add_rect(slide, Inches(0.8), Inches(1.05), Inches(2.0), Inches(0.04), ACCENT)

    add_text(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.6),
             "カットリストに登録したセグメントを1つの動画ファイルとして出力します。",
             font_size=18, color=WHITE)

    # エクスポートモーダル風カード
    modal_x, modal_y = Inches(2.5), Inches(2.3)
    modal_w, modal_h = Inches(8), Inches(4.5)
    add_rect(slide, modal_x, modal_y, modal_w, modal_h, BG_SURFACE, corner_radius=0.04)
    add_rect(slide, modal_x, modal_y, modal_w, Inches(0.06), GREEN)

    add_text(slide, modal_x + Inches(0.5), modal_y + Inches(0.3), Inches(3), Inches(0.5),
             "エクスポート", font_size=20, color=WHITE, bold=True)

    info_items = [
        "ソース:  sample.mp4",
        "セグメント数:  2",
        "合計時間:  00:01:05.00",
        "上下反転:  OFF",
        "出力:  無劣化コピー（再エンコードなし）",
        "保存先:  ソースファイルと同じフォルダ",
    ]
    for i, item in enumerate(info_items):
        add_text(slide, modal_x + Inches(0.5), modal_y + Inches(1.0) + i * Inches(0.4),
                 Inches(7), Inches(0.4), item, font_size=14, color=GRAY)

    # ボタン
    shape = add_rect(slide, modal_x + modal_w - Inches(2.6), modal_y + modal_h - Inches(0.8),
                     Inches(1.0), Inches(0.45), BG_CARD, corner_radius=0.1)
    shape.line.color.rgb = RGBColor(0x44, 0x44, 0x44)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.paragraphs[0].text = "キャンセル"
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    shape = add_rect(slide, modal_x + modal_w - Inches(1.4), modal_y + modal_h - Inches(0.8),
                     Inches(1.2), Inches(0.45), RGBColor(0x1B, 0x4D, 0x1E), corner_radius=0.1)
    shape.line.color.rgb = GREEN
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.paragraphs[0].text = "エクスポート"
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = GREEN
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 出力ファイル名の説明
    add_rect(slide, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.8), BG_SURFACE, corner_radius=0.04)
    add_text(slide, Inches(1.1), Inches(6.3), Inches(11), Inches(0.6),
             "出力ファイル名:  元ファイル名_cut.mp4  （重複時は _cut_1, _cut_2 ... と連番）\n"
             "保存場所:  元の動画ファイルと同じフォルダ",
             font_size=14, color=GRAY)

    # --- 9. トラブルシューティング ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "トラブルシューティング", font_size=32, color=ACCENT, bold=True)
    add_rect(slide, Inches(0.8), Inches(1.05), Inches(3.6), Inches(0.04), ACCENT)

    troubles = [
        ("「開発元を確認できない」と表示される",
         "初回のみ、アプリを右クリック →「開く」を選択してください。\n"
         "2回目以降は通常どおりダブルクリックで起動できます。"),
        ("「ffmpeg が必要です」と表示される",
         "ターミナルを開いて brew install ffmpeg を実行してください。\n"
         "Homebrew 未導入の場合は先に brew.sh からインストールしてください。"),
        ("動画が再生されない / シークできない",
         "ブラウザが対応していないコーデックの可能性があります。\n"
         "MP4（H.264）形式の動画をお試しください。"),
        ("エクスポートに時間がかかる",
         "上下反転ONの場合は再エンコードが必要なため時間がかかります。\n"
         "反転不要であればOFFにすると無劣化コピーで高速に処理されます。"),
    ]
    for i, (q, a) in enumerate(troubles):
        y = Inches(1.5) + i * Inches(1.45)
        add_rect(slide, Inches(0.8), y, Inches(11.7), Inches(1.25), BG_CARD, corner_radius=0.04)
        add_text(slide, Inches(1.1), y + Inches(0.1), Inches(11), Inches(0.4),
                 "Q: " + q, font_size=15, color=RED, bold=True)
        add_text(slide, Inches(1.1), y + Inches(0.55), Inches(11), Inches(0.6),
                 a, font_size=13, color=GRAY)

    # --- 10. 最終スライド ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    card_w, card_h = Inches(8), Inches(3.5)
    card_x = (SLIDE_W - card_w) // 2
    card_y = (SLIDE_H - card_h) // 2 - Inches(0.3)
    add_rect(slide, card_x, card_y, card_w, card_h, BG_SURFACE, corner_radius=0.05)
    add_rect(slide, card_x, card_y, card_w, Inches(0.06), ACCENT)

    add_text(slide, card_x, card_y + Inches(0.6), card_w, Inches(0.8),
             "VideoCut", font_size=44, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, card_x, card_y + Inches(1.5), card_w, Inches(0.5),
             "シンプルに、素早く、動画をカット。", font_size=20, color=WHITE,
             alignment=PP_ALIGN.CENTER)
    add_text(slide, card_x, card_y + Inches(2.4), card_w, Inches(0.5),
             "github.com/IvyTechnx/VideoCut", font_size=14, color=GRAY,
             alignment=PP_ALIGN.CENTER)

    # --- 保存 ---
    prs.save(OUT)
    print(f"Manual saved: {OUT}")


if __name__ == "__main__":
    build()
